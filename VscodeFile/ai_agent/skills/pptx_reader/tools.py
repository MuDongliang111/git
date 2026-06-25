"""提供给 pptx_reader 技能的自定义工具集。

包含以下工具：
- read_pptx: 完整读取 PPTX 文件，提取文本、表格、备注
- search_pptx: 在 PPTX 文件中搜索指定关键词
- get_pptx_outline: 提取幻灯片大纲（仅标题和结构）
- get_pptx_metadata: 获取文档元数据（作者、创建日期、幻灯片数等）
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from ai_agent.tools.base import ToolDefinition


# =============================================================================
# 通用辅助函数
# =============================================================================


def _resolve_path(path: str) -> Path:
    """解析并验证 .pptx 文件路径。"""
    file_path = Path(path)
    if not file_path.is_absolute():
        file_path = Path.cwd() / file_path
    return file_path


def _validate_pptx(file_path: Path) -> str | None:
    """验证文件是否为有效的 .pptx 文件。返回错误信息或 None。"""
    if not file_path.exists():
        return f"错误: 文件不存在: {file_path}"
    if not file_path.suffix.lower() == ".pptx":
        return f"错误: 不支持的文件格式 '{file_path.suffix}'，仅支持 .pptx 文件"
    return None


def _ensure_pptx():
    """确保 python-pptx 库可用。"""
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        raise ImportError(
            "错误: 未安装 python-pptx 库。请运行以下命令安装:\n"
            "  pip install python-pptx"
        )


def _extract_shape_text(shape) -> list[str]:
    """从形状中提取所有文本段落。"""
    texts: list[str] = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                texts.append(text)
    return texts


def _extract_table_text(shape) -> list[str]:
    """从形状中提取表格内容。"""
    rows: list[str] = []
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            row_cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(row_cells))
    return rows


def _format_slide_content(
    slide,
    slide_num: int,
    include_notes: bool = True,
) -> str:
    """格式化单张幻灯片的内容。"""
    lines: list[str] = []
    lines.append(f"\n## 幻灯片 {slide_num}")

    # 标题
    title_texts: list[str] = []
    if slide.shapes.title and slide.shapes.title.text.strip():
        title_texts.append(slide.shapes.title.text.strip())

    # 正文和表格
    body_texts: list[str] = []
    table_texts: list[str] = []

    for shape in slide.shapes:
        # 跳过已处理的标题形状
        if shape == slide.shapes.title:
            continue

        body_texts.extend(_extract_shape_text(shape))
        table_texts.extend(_extract_table_text(shape))

    # 输出标题
    if title_texts:
        for t in title_texts:
            lines.append(f"  标题: {t}")

    # 输出正文
    if body_texts:
        lines.append("  正文:")
        for t in body_texts:
            lines.append(f"    • {t}")

    # 输出表格
    if table_texts:
        lines.append("  表格:")
        for t in table_texts:
            lines.append(f"    | {t} |")

    # 备注
    if include_notes:
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    lines.append(f"  备注: {notes_text}")
        except Exception:
            pass  # 某些情况下备注可能无法访问

    # 无内容提示
    if not title_texts and not body_texts and not table_texts:
        lines.append("  （无文本内容）")

    return "\n".join(lines)


# =============================================================================
# 工具处理函数
# =============================================================================


def _read_pptx(
    path: str,
    include_notes: bool = True,
    max_slides: int = 0,
    slide_numbers: str = "",
) -> str:
    """读取 .pptx 文件并提取所有幻灯片的文本内容。

    Parameters
    ----------
    path : str
        要读取的 .pptx 文件路径。
    include_notes : bool
        是否包含演讲者备注。
    max_slides : int
        最大读取幻灯片数（0 表示不限制）。
    slide_numbers : str
        要读取的特定幻灯片编号，如 "1,3,5-7"。
    """
    _ensure_pptx()
    from pptx import Presentation

    file_path = _resolve_path(path)
    error = _validate_pptx(file_path)
    if error:
        return error

    try:
        prs = Presentation(str(file_path))
    except Exception as exc:
        return f"错误: 无法打开 PowerPoint 文件 '{file_path}': {exc}"

    total_slides = len(prs.slides)

    # 解析要读取的幻灯片范围
    target_slides = _parse_slide_range(slide_numbers, total_slides)
    if target_slides is not None:
        # 使用指定的幻灯片编号
        slide_set = set(target_slides)
    else:
        # 使用 max_slides 限制
        limit = max_slides if max_slides > 0 else total_slides
        limit = min(limit, total_slides)
        slide_set = set(range(1, limit + 1))

    # 构建输出
    lines: list[str] = []
    lines.append(f"文件: {file_path.name}")
    lines.append(f"幻灯片总数: {total_slides}")
    if slide_numbers:
        lines.append(f"（已筛选幻灯片: {slide_numbers}）")
    elif max_slides > 0 and max_slides < total_slides:
        lines.append(f"（截取前 {max_slides} 张）")
    lines.append("=" * 50)

    for slide_num, slide in enumerate(prs.slides, start=1):
        if slide_num not in slide_set:
            continue
        lines.append(
            _format_slide_content(slide, slide_num, include_notes=include_notes)
        )

    total_read = sum(1 for s in slide_set if 1 <= s <= total_slides)
    if total_read < len(slide_set):
        lines.append(
            f"\n注意: 请求了 {len(slide_set)} 张幻灯片，"
            f"但文件只有 {total_slides} 张"
        )

    return "\n".join(lines)


def _search_pptx(
    path: str,
    query: str,
    case_sensitive: bool = False,
    include_notes: bool = True,
    max_results: int = 50,
) -> str:
    """在 PPTX 文件中搜索关键词，返回匹配的幻灯片和上下文。

    Parameters
    ----------
    path : str
        要搜索的 .pptx 文件路径。
    query : str
        搜索关键词。
    case_sensitive : bool
        是否区分大小写。
    include_notes : bool
        是否同时搜索演讲者备注。
    max_results : int
        最大返回结果数。
    """
    _ensure_pptx()
    from pptx import Presentation

    file_path = _resolve_path(path)
    error = _validate_pptx(file_path)
    if error:
        return error

    try:
        prs = Presentation(str(file_path))
    except Exception as exc:
        return f"错误: 无法打开 PowerPoint 文件 '{file_path}': {exc}"

    search_query = query if case_sensitive else query.lower()
    results: list[dict[str, Any]] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_matches: list[dict[str, str]] = []

        # 搜索标题
        if slide.shapes.title and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()
            cmp_title = title if case_sensitive else title.lower()
            if search_query in cmp_title:
                slide_matches.append({
                    "location": "标题",
                    "text": _truncate_context(title, query, case_sensitive),
                })

        # 搜索正文
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    cmp_text = text if case_sensitive else text.lower()
                    if search_query in cmp_text:
                        slide_matches.append({
                            "location": "正文",
                            "text": _truncate_context(text, query, case_sensitive),
                        })
                        if len(slide_matches) >= 10:
                            break

            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    cmp_row = row_text if case_sensitive else row_text.lower()
                    if search_query in cmp_row:
                        slide_matches.append({
                            "location": "表格",
                            "text": _truncate_context(row_text, query, case_sensitive),
                        })

        # 搜索备注
        if include_notes:
            try:
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    cmp_notes = notes_text if case_sensitive else notes_text.lower()
                    if search_query in cmp_notes:
                        slide_matches.append({
                            "location": "备注",
                            "text": _truncate_context(notes_text, query, case_sensitive),
                        })
            except Exception:
                pass

        if slide_matches:
            results.append({
                "slide": slide_num,
                "count": len(slide_matches),
                "matches": slide_matches[:5],  # 每张幻灯片最多展示 5 条
            })

        if len(results) >= max_results:
            break

    if not results:
        return (
            f"在文件 '{file_path.name}' 中未找到 '{query}' 的相关内容。\n"
            f"（共搜索了 {len(prs.slides)} 张幻灯片）"
        )

    total_matches = sum(r["count"] for r in results)
    lines: list[str] = []
    lines.append(
        f"搜索 '{query}' — 在 {file_path.name} 中找到 "
        f"{total_matches} 条匹配（{len(results)} 张幻灯片）"
    )
    lines.append("=" * 50)

    for result in results:
        lines.append(f"\n## 幻灯片 {result['slide']} ({result['count']} 条匹配)")
        for i, match in enumerate(result["matches"], 1):
            lines.append(f"  [{match['location']}] {match['text']}")

    if total_matches > sum(len(r["matches"]) for r in results):
        lines.append("\n注意: 部分匹配结果已截断，请缩小搜索范围获取更精确的结果。")

    return "\n".join(lines)


def _get_pptx_outline(path: str) -> str:
    """提取 PPTX 文件的大纲结构（仅标题和章节）。

    Parameters
    ----------
    path : str
        要读取的 .pptx 文件路径。
    """
    _ensure_pptx()
    from pptx import Presentation

    file_path = _resolve_path(path)
    error = _validate_pptx(file_path)
    if error:
        return error

    try:
        prs = Presentation(str(file_path))
    except Exception as exc:
        return f"错误: 无法打开 PowerPoint 文件 '{file_path}': {exc}"

    total_slides = len(prs.slides)
    lines: list[str] = []
    lines.append(f"文件: {file_path.name}")
    lines.append(f"幻灯片总数: {total_slides}")
    lines.append(f"幻灯片尺寸: {prs.slide_width} x {prs.slide_height} (EMU)")
    lines.append("=" * 50)

    for slide_num, slide in enumerate(prs.slides, start=1):
        # 获取标题
        title = ""
        if slide.shapes.title and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()

        # 统计内容
        text_shapes = 0
        table_shapes = 0
        image_shapes = 0
        chart_shapes = 0
        total_text_length = 0

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                text_shapes += 1
                for para in shape.text_frame.paragraphs:
                    total_text_length += len(para.text.strip())
            if shape.has_table:
                table_shapes += 1
            if shape.shape_type is not None:
                # MSO_SHAPE_TYPE: 13 = Picture
                if getattr(shape.shape_type, "__name__", "") == "PICTURE" or str(shape.shape_type) == "PICTURE (13)":
                    image_shapes += 1
                # MSO_SHAPE_TYPE: 3 = Chart
                if hasattr(shape, "has_chart") and shape.has_chart:
                    chart_shapes += 1

        # 布局名称
        layout_name = ""
        try:
            if slide.slide_layout:
                layout_name = slide.slide_layout.name
        except Exception:
            pass

        # 格式化输出
        if title:
            lines.append(f"\n## 幻灯片 {slide_num}: {title}")
        else:
            lines.append(f"\n## 幻灯片 {slide_num}: （无标题）")

        details: list[str] = []
        if layout_name:
            details.append(f"布局: {layout_name}")
        if text_shapes:
            details.append(f"文本框: {text_shapes} 个 ({total_text_length} 字符)")
        if table_shapes:
            details.append(f"表格: {table_shapes} 个")
        if image_shapes:
            details.append(f"图片: {image_shapes} 张")
        if chart_shapes:
            details.append(f"图表: {chart_shapes} 个")

        if details:
            lines.append(f"  ({', '.join(details)})")

        # 尝试获取 section 信息
        try:
            # python-pptx 没有直接的 section API，通过 XML 解析
            if hasattr(slide, "_element"):
                el = slide._element
                # 查找前一个 section 分隔符（如果有的话）
                prev = el.getprevious()
                if prev is not None:
                    sld_id = el.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
                    # 简化的 section 检测（不完美，但可用）
        except Exception:
            pass

    # 统计信息
    total_text_shapes = 0
    total_tables = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape != slide.shapes.title and shape.has_text_frame:
                total_text_shapes += 1
            if shape.has_table:
                total_tables += 1

    lines.append(f"\n{'=' * 50}")
    lines.append(f"摘要: {total_slides} 张幻灯片, {total_text_shapes} 个文本框, {total_tables} 个表格")

    return "\n".join(lines)


def _get_pptx_metadata(path: str) -> str:
    """获取 PPTX 文件的元数据信息。

    Parameters
    ----------
    path : str
        要查看的 .pptx 文件路径。
    """
    _ensure_pptx()
    from pptx import Presentation

    file_path = _resolve_path(path)
    error = _validate_pptx(file_path)
    if error:
        return error

    # 文件系统信息
    stat = file_path.stat()
    file_size_kb = stat.st_size / 1024

    try:
        prs = Presentation(str(file_path))
    except Exception as exc:
        return f"错误: 无法打开 PowerPoint 文件 '{file_path}': {exc}"

    # 核心属性
    core_props = prs.core_properties

    lines: list[str] = []
    lines.append(f"文件信息: {file_path.name}")
    lines.append("=" * 50)
    lines.append(f"文件路径:     {file_path}")
    lines.append(f"文件大小:     {file_size_kb:.1f} KB")
    lines.append(f"幻灯片数:     {len(prs.slides)}")
    lines.append(f"幻灯片宽度:   {prs.slide_width} EMU ({prs.slide_width / 914400:.1f} 英寸)")
    lines.append(f"幻灯片高度:   {prs.slide_height} EMU ({prs.slide_height / 914400:.1f} 英寸)")

    # 文档属性
    lines.append(f"\n文档属性:")
    lines.append("-" * 30)

    prop_fields = [
        ("标题", core_props.title),
        ("作者", core_props.author),
        ("主题", core_props.subject),
        ("关键词", core_props.keywords),
        ("类别", core_props.category),
        ("描述", core_props.comments),
        ("状态", core_props.content_status),
        ("修订版本", core_props.revision),
        ("创建者", core_props.last_modified_by),
        ("创建时间", str(core_props.created) if core_props.created else None),
        ("修改时间", str(core_props.modified) if core_props.modified else None),
        ("上次打印", str(core_props.last_printed) if core_props.last_printed else None),
        ("语言", core_props.language),
        ("标识符", core_props.identifier),
    ]

    for label, value in prop_fields:
        if value:
            lines.append(f"  {label}: {value}")

    # 幻灯片布局统计
    layout_counts: dict[str, int] = {}
    for slide in prs.slides:
        try:
            layout_name = slide.slide_layout.name if slide.slide_layout else "未知"
        except Exception:
            layout_name = "未知"
        layout_counts[layout_name] = layout_counts.get(layout_name, 0) + 1

    if len(layout_counts) > 1:
        lines.append(f"\n幻灯片布局分布:")
        lines.append("-" * 30)
        for layout_name, count in sorted(layout_counts.items(), key=lambda x: -x[1]):
            lines.append(f"  {layout_name}: {count} 张")

    # 内容统计
    total_text = 0
    total_images = 0
    total_tables = 0
    total_charts = 0
    total_notes = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    total_text += len(para.text)
            if shape.has_table:
                total_tables += 1
            if hasattr(shape, "has_chart") and shape.has_chart:
                total_charts += 1
            # 检测图片
            shape_type_str = str(shape.shape_type) if shape.shape_type else ""
            if "PICTURE" in shape_type_str or "Picture" in shape_type_str:
                total_images += 1
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    total_notes += 1
        except Exception:
            pass

    lines.append(f"\n内容统计:")
    lines.append("-" * 30)
    lines.append(f"  总字符数:   {total_text:,}")
    lines.append(f"  图片:       {total_images} 张")
    lines.append(f"  表格:       {total_tables} 个")
    lines.append(f"  图表:       {total_charts} 个")
    lines.append(f"  含备注幻灯片: {total_notes} 张")

    # 幻灯片母版
    try:
        num_masters = len(prs.slide_masters) if prs.slide_masters else 0
        if num_masters:
            lines.append(f"\n母版信息:")
            lines.append("-" * 30)
            lines.append(f"  幻灯片母版: {num_masters} 个")
            for i, master in enumerate(prs.slide_masters, 1):
                layouts = len(master.slide_layouts) if master.slide_layouts else 0
                lines.append(f"  母版 {i}: {layouts} 种布局")
    except Exception:
        pass

    return "\n".join(lines)


# =============================================================================
# 内部辅助函数
# =============================================================================


def _parse_slide_range(spec: str, total: int) -> list[int] | None:
    """解析幻灯片范围字符串。

    示例: "1,3,5-7" -> [1, 3, 5, 6, 7]
    返回 None 表示使用默认行为。
    """
    if not spec or not spec.strip():
        return None

    result: list[int] = []
    parts = re.split(r"[,;，；]\s*", spec.strip())
    for part in parts:
        part = part.strip()
        if "-" in part:
            try:
                start_s, end_s = part.split("-", 1)
                start = int(start_s.strip())
                end = int(end_s.strip())
                if start > end:
                    start, end = end, start
                result.extend(range(start, end + 1))
            except ValueError:
                continue
        else:
            try:
                result.append(int(part))
            except ValueError:
                continue

    # 去重并限制在有效范围内
    seen: set[int] = set()
    valid: list[int] = []
    for n in result:
        if 1 <= n <= total and n not in seen:
            seen.add(n)
            valid.append(n)

    return sorted(valid) if valid else None


def _truncate_context(text: str, query: str, case_sensitive: bool) -> str:
    """截取匹配关键词周围的上下文文本。"""
    max_len = 120
    if len(text) <= max_len:
        return text

    # 找到匹配位置
    search_query = query if case_sensitive else query.lower()
    search_text = text if case_sensitive else text.lower()
    idx = search_text.find(search_query)

    if idx < 0:
        return text[:max_len] + "..."

    # 截取上下文
    start = max(0, idx - 40)
    end = min(len(text), idx + len(query) + 60)
    result = text[start:end]
    if start > 0:
        result = "..." + result
    if end < len(text):
        result = result + "..."
    return result


# =============================================================================
# JSON Schema 定义
# =============================================================================

_READ_PPTX_SCHEMA = {
    "type": "object",
    "properties": {
        "path": {
            "type": "string",
            "description": "要读取的 .pptx 文件路径（相对或绝对路径）",
        },
        "include_notes": {
            "type": "boolean",
            "description": "是否包含演讲者备注（默认 true）",
            "default": True,
        },
        "max_slides": {
            "type": "integer",
            "description": "最大读取幻灯片数（默认 0 表示不限制）",
            "default": 0,
        },
        "slide_numbers": {
            "type": "string",
            "description": (
                "要读取的特定幻灯片编号，如 '1,3,5-7'。"
                "留空表示读取所有幻灯片或按 max_slides 限制"
            ),
            "default": "",
        },
    },
    "required": ["path"],
}

_SEARCH_PPTX_SCHEMA = {
    "type": "object",
    "properties": {
        "path": {
            "type": "string",
            "description": "要搜索的 .pptx 文件路径（相对或绝对路径）",
        },
        "query": {
            "type": "string",
            "description": "搜索关键词",
        },
        "case_sensitive": {
            "type": "boolean",
            "description": "是否区分大小写（默认 false）",
            "default": False,
        },
        "include_notes": {
            "type": "boolean",
            "description": "是否同时搜索演讲者备注（默认 true）",
            "default": True,
        },
        "max_results": {
            "type": "integer",
            "description": "最大返回的幻灯片结果数（默认 50）",
            "default": 50,
        },
    },
    "required": ["path", "query"],
}

_OUTLINE_PPTX_SCHEMA = {
    "type": "object",
    "properties": {
        "path": {
            "type": "string",
            "description": "要提取大纲的 .pptx 文件路径（相对或绝对路径）",
        },
    },
    "required": ["path"],
}

_METADATA_PPTX_SCHEMA = {
    "type": "object",
    "properties": {
        "path": {
            "type": "string",
            "description": "要查看元数据的 .pptx 文件路径（相对或绝对路径）",
        },
    },
    "required": ["path"],
}


# =============================================================================
# 工具注册入口
# =============================================================================


def get_tools() -> list[ToolDefinition]:
    """返回此技能提供的工具列表。"""
    return [
        ToolDefinition(
            name="read_pptx",
            description=(
                "读取 PowerPoint (.pptx) 文件，提取所有幻灯片的文本内容，"
                "包括标题、正文、表格文本和演讲者备注。"
                "支持筛选特定幻灯片编号（如 '1,3,5-7'）和限制最大幻灯片数。"
                "适用于查看、总结和分析 PPT 内容。"
            ),
            parameters=_READ_PPTX_SCHEMA,
            handler=_read_pptx,
            source="skill",
            source_name="pptx_reader",
        ),
        ToolDefinition(
            name="search_pptx",
            description=(
                "在 PowerPoint (.pptx) 文件中搜索指定关键词，"
                "返回匹配的幻灯片编号、位置（标题/正文/表格/备注）和上下文文本。"
                "支持区分大小写搜索和最大结果数限制。"
                "适用于快速定位 PPT 中的特定内容。"
            ),
            parameters=_SEARCH_PPTX_SCHEMA,
            handler=_search_pptx,
            source="skill",
            source_name="pptx_reader",
        ),
        ToolDefinition(
            name="get_pptx_outline",
            description=(
                "提取 PowerPoint (.pptx) 文件的完整大纲结构，"
                "包括每张幻灯片的标题、布局类型和内容统计"
                "（文本框数、表格数、图片数、图表数）。"
                "适用于快速了解 PPT 的整体结构和组织方式。"
            ),
            parameters=_OUTLINE_PPTX_SCHEMA,
            handler=_get_pptx_outline,
            source="skill",
            source_name="pptx_reader",
        ),
        ToolDefinition(
            name="get_pptx_metadata",
            description=(
                "获取 PowerPoint (.pptx) 文件的详细元数据信息，"
                "包括文件大小、幻灯片尺寸、文档属性（作者、标题、创建/修改时间等）、"
                "布局分布、内容统计（字符数、图片、表格、图表、备注数）和母版信息。"
                "适用于全面了解 PPT 文件的技术属性和文档信息。"
            ),
            parameters=_METADATA_PPTX_SCHEMA,
            handler=_get_pptx_metadata,
            source="skill",
            source_name="pptx_reader",
        ),
    ]
